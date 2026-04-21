from __future__ import annotations

import argparse
import xml.etree.ElementTree as ET
from pathlib import Path
from collections.abc import Iterable

from markitdown import MarkItDown
from pptx import Presentation


SKIP_EXTENSIONS = {".md"}
PPTX_SUFFIX = ".pptx"
SMARTART_GRAPHIC_DATA_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
SMARTART_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"


def iter_input_files(data_dir: Path) -> Iterable[Path]:
    """Yield all files under data_dir recursively."""
    for path in data_dir.rglob("*"):
        if path.is_file() and path.suffix.lower() not in SKIP_EXTENSIONS:
            yield path


def _collect_text_nodes(xml_root: ET.Element) -> list[str]:
    """Collect every text node from an XML element tree."""
    texts: list[str] = []

    for node in xml_root.iter():
        if node.tag.rsplit("}", 1)[-1] != "t":
            continue

        text = (node.text or "").strip()
        if text:
            texts.append(text)

    return texts


def _normalize_text_block(text: str) -> str:
    """Normalize a SmartArt text block for de-duplication."""
    # Splitlines will split by markdown spaces we preserve (e.g. `  * `)
    # We should NOT join the lines into one line anymore because Markdown lists require line breaks!
    lines = [line.rstrip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _is_smartart_shape(shape) -> bool:
    """Return True when a shape is a SmartArt graphic frame."""
    if not hasattr(shape, "element") or shape.element.tag.rsplit("}", 1)[-1] != "graphicFrame":
        return False

    try:
        return shape.element.graphic.graphicData.uri == SMARTART_GRAPHIC_DATA_URI
    except AttributeError:
        return False


def get_smartart_elements(shape) -> list[dict[str, object]]:
    elements = []

    if not hasattr(shape, "part"):
        return elements

    namespaces = {
        "dsp": "http://schemas.microsoft.com/office/drawing/2008/diagram",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    try:
        drawing_parts = [r.target_part for r in shape.part.rels.values() if r.reltype.endswith("/diagramDrawing")]
        if not drawing_parts: return elements
        drawing_root = ET.fromstring(drawing_parts[0].blob)
    except Exception:
        return elements

    # Build a parent map to find nearest a:off
    parent_map = {c: p for p in drawing_root.iter() for c in p}

    def get_nearest_offset(element) -> tuple[int, int]:
        curr = element
        while curr is not None:
            # xfrm that directly holds off
            xfrm = curr.find(".//a:xfrm", namespaces)
            if xfrm is not None:
                off = xfrm.find(".//a:off", namespaces)
                if off is not None:
                    return int(off.get("x", 0)), int(off.get("y", 0))

            # Sometimes it's inside spPr
            spPr = curr.find(".//dsp:spPr/a:xfrm", namespaces)
            if spPr is not None:
                off = spPr.find(".//a:off", namespaces)
                if off is not None:
                    return int(off.get("x", 0)), int(off.get("y", 0))
            
            curr = parent_map.get(curr)
        return 0, 0

    # Universal Extraction: directly find all paragraph tags globally
    # to avoid omitting any text node block.
    for p in drawing_root.findall(".//a:p", namespaces):
        run_texts = [t.text for t in p.findall(".//a:t", namespaces) if t.text]
        block_text = "".join(run_texts).strip()

        if block_text:
            x, y = get_nearest_offset(p)
            elements.append({
                "text": block_text,
                "left": getattr(shape, "left", 0) + x,
                "top": getattr(shape, "top", 0) + y
            })

    return elements


import re

def get_text_frame_text(shape) -> str:
    """Extract normalized text from a regular text frame with Markdown formatting."""
    if not getattr(shape, "has_text_frame", False):
        return ""

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return ""

    lines = []
    for p in text_frame.paragraphs:
        p_text = ""
        for run in p.runs:
            run_text = run.text.replace("\n", " ")
            if getattr(run.font, "bold", False) and run_text.strip():
                stripped = run_text.strip()
                # Wrap with markdown bold unless it's already wrapped (basic check)
                run_text = run_text.replace(stripped, f"**{stripped}**")
            p_text += run_text
        
        p_text = p_text.strip()
        if not p_text:
            continue
            
        # Strip existing PPT list bullet characters
        p_text = re.sub(r'^([o\uf0b7\u25a0\u25cb\u2022\u25cf\uf0d8\u27a2\u25aa\u25a1]|\?)\s+', '', p_text)
        
        # Heading Detection
        is_heading = False
        # 1. Matches patterns like 5.1.4
        if re.match(r'^\d+(\.\d+)+\s+', p_text):
            is_heading = True
        # 2. Heuristic for font size (if available and large enough, though runs can vary, we assume 18+ is heading)
        elif p.runs and p.runs[0].font and p.runs[0].font.size and p.runs[0].font.size.pt > 20:
            is_heading = True
            
        level = getattr(p, "level", 0)
        
        if is_heading:
            # Decide header level (e.g. 5.1 -> ##, 5.1.4 -> ###)
            h_match = re.match(r'^(\d+(?:\.\d+)+)', p_text)
            if h_match:
                dots = h_match.group(1).count('.')
                h_prefix = "#" * min(dots + 2, 4)
                lines.append(f"{h_prefix} {p_text}")
            else:
                lines.append(f"## {p_text}")
        else:
            indent = "  " * level
            lines.append(f"{indent}* {p_text}")

    return "\n".join(lines)


def reconstruct_table(elements: list[dict]) -> str | None:
    """Attempt to reconstruct a Markdown table from scattered text blocks."""
    if not elements:
        return None

    # Group elements by Top (Row) with 10% tolerance (using average top gap)
    # Actually, a simple absolute tolerance is usually safer for PPTX EMUs.
    # Let's say 100,000 EMU vertical tolerance for items in the same row.
    ROW_TOLERANCE = 150000 
    rows_dict = []
    
    # Sort elements by top to process sequentially
    elements_sorted = sorted(elements, key=lambda x: x["top"])
    
    for el in elements_sorted:
        placed = False
        for row in rows_dict:
            if abs(el["top"] - row["top_avg"]) < ROW_TOLERANCE:
                row["elements"].append(el)
                row["top_avg"] = sum(x["top"] for x in row["elements"]) / len(row["elements"])
                placed = True
                break
        if not placed:
            rows_dict.append({"top_avg": el["top"], "elements": [el]})

    # Find column baselines by aggregating Left coordinates
    left_coords = [el["left"] for el in elements]
    COL_TOLERANCE = 300000 # Horizontal tolerance to cluster lefts
    
    col_clusters = []
    for left in sorted(left_coords):
        placed = False
        for cluster in col_clusters:
            if abs(left - cluster["avg"]) < COL_TOLERANCE:
                cluster["count"] += 1
                cluster["avg"] = ((cluster["avg"] * (cluster["count"] - 1)) + left) / cluster["count"]
                placed = True
                break
        if not placed:
            col_clusters.append({"avg": left, "count": 1})
            
    # Filter valid column baselines - must pop up frequently enough (e.g. at least 2 items)
    valid_cols = sorted([c["avg"] for c in col_clusters if c["count"] >= 2])
    
    # Needs at least 2 columns to be considered a table
    if len(valid_cols) < 2:
        return None

    lines = []
    for row in rows_dict:
        row_elements = row["elements"]
        
        # If it's a single element spanning, treat it as a special block (not table)
        if len(row_elements) == 1:
            lines.append(f"\n{row_elements[0]['text']}")
            continue
            
        # Map elements to columns
        row_cells = [""] * len(valid_cols)
        for el in row_elements:
            # Find closest valid col
            closest_col_idx = min(range(len(valid_cols)), key=lambda i: abs(valid_cols[i] - el["left"]))
            # If multiple elements map to the same cell, join them safely
            if row_cells[closest_col_idx]:
                row_cells[closest_col_idx] += " " + el["text"].replace("\n", " ")
            else:
                row_cells[closest_col_idx] = el["text"].replace("\n", " ")
        
        # Build Markdown table row
        lines.append("| " + " | ".join(row_cells) + " |")
        
        # Add separator after first table row
        if len(lines) == 1 or (len(lines) > 1 and not lines[-2].strip().startswith("|")):
            lines.append("|" + "|".join(["---"] * len(valid_cols)) + "|")

    return "\n".join(lines)


def slide_to_markdown(slide, slide_number: int) -> str:
    """Render a single slide into markdown using grid-based spatial ordering."""
    elements_on_slide: list[dict[str, object]] = []

    def process_shape(shape, group_left=None, group_top=None):
        left = group_left if group_left is not None else getattr(shape, "left", 0)
        
        if getattr(shape, "shape_type", None) == 6:  # GROUP
            g_left = getattr(shape, "left", 0)
            g_top = getattr(shape, "top", 0)
            for i, sub_shape in enumerate(shape.shapes):
                process_shape(sub_shape, group_left=g_left, group_top=g_top + (i * 1000))
            return

        if getattr(shape, "has_text_frame", False):
            text = get_text_frame_text(shape)
            if text:
                normalized_text = _normalize_text_block(text)
                if normalized_text:
                    top_val = group_top if group_top is not None else getattr(shape, "top", 0)
                    elements_on_slide.append({"text": normalized_text, "top": top_val, "left": left})
            return

        if _is_smartart_shape(shape):
            for sa_item in get_smartart_elements(shape):
                normalized_block = _normalize_text_block(sa_item["text"])
                if normalized_block:
                    elements_on_slide.append({"text": normalized_block, "top": sa_item["top"], "left": sa_item["left"]})

    for shape in slide.shapes:
        process_shape(shape)

    # Precise De-duplication (only remove if text matches AND coordinates are extremely close ~1000 EMU)
    deduped_elements: list[dict[str, object]] = []
    seen_objects: list[dict[str, object]] = []
    
    for el in elements_on_slide:
        is_dupe = False
        for seen in seen_objects:
            if el["text"] == seen["text"] and abs(int(el["left"]) - int(seen["left"])) < 1000 and abs(int(el["top"]) - int(seen["top"])) < 1000:
                is_dupe = True
                break
        if not is_dupe:
            seen_objects.append(el)
            deduped_elements.append(el)

    # Segment the slide vertically into Header, Body, and Footer
    headers = []
    body = []
    footers = []

    for el in deduped_elements:
        t = int(el["top"])
        if t < 1000000:
            headers.append(el)
        elif t > 5000000:
            footers.append(el)
        else:
            body.append(el)

    headers.sort(key=lambda item: item["top"])
    footers.sort(key=lambda item: (item["top"], item["left"]))

    # Primary sort body by raw Left then Top
    body.sort(key=lambda item: (item["left"], item["top"]))

    # Table Reconstruction / Grid-Based Column grouping for Body
    is_table = False
    
    # Simple check: do we have many items sharing the same top coordinates?
    if len(body) > 4:
        row_counts = {}
        for b in body:
            top_rounded = (int(b["top"]) // 100000) * 100000
            row_counts[top_rounded] = row_counts.get(top_rounded, 0) + 1
        
        # If multiple rows have >1 item, it's very likely a table.
        if sum(1 for count in row_counts.values() if count > 1) >= 2:
            is_table = True

    final_elements = []
    if is_table:
        table_md = reconstruct_table(body)
        if table_md:
            final_elements = headers + [{"text": table_md, "top": 0, "left": 0}] + footers
        else:
            is_table = False  # Fallback

    if not is_table:
        # Flowchart logic: group into Columns
        TOLERANCE = 600000
        columns: list[list[dict[str, object]]] = []
        
        for el in body:
            placed = False
            for col in columns:
                if abs(int(el["left"]) - int(col[0]["left"])) < TOLERANCE:
                    col.append(el)
                    placed = True
                    break
            if not placed:
                columns.append([el])

        columns.sort(key=lambda col: sum(int(item["left"]) for item in col) / len(col))

        for col in columns:
            col.sort(key=lambda item: item["top"])

        final_elements = headers + [item for col in columns for item in col] + footers

    lines = [f"## Slide {slide_number}", ""]
    for item in final_elements:
        lines.append(str(item["text"]).strip())
        lines.append("")

    return "\n".join(lines).rstrip()


def extract_all_content(pptx_path: Path, md_converter: MarkItDown | None = None) -> str:
    """Convert a PPTX into slide-ordered markdown with inline SmartArt text."""
    del md_converter

    presentation = Presentation(str(pptx_path))
    slide_blocks: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_blocks.append(slide_to_markdown(slide, slide_number))

    return "\n\n".join(block for block in slide_blocks if block).strip()


def convert_file(md_converter: MarkItDown, src: Path, dst: Path) -> None:
    """Convert one source file to markdown and write to dst."""
    if src.suffix.lower() == PPTX_SUFFIX:
        markdown = extract_all_content(src, md_converter)
    else:
        result = md_converter.convert(str(src))
        markdown = result.text_content

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(markdown, encoding="utf-8")


def process_all(data_dir: Path, output_dir: Path) -> tuple[int, int]:
    """Process all files and return (success_count, fail_count)."""
    md_converter = MarkItDown()
    success_count = 0
    fail_count = 0

    for src in iter_input_files(data_dir):
        relative = src.relative_to(data_dir)
        dst = output_dir / relative.with_suffix(".md")

        try:
            convert_file(md_converter, src, dst)
            success_count += 1
            print(f"[OK]   {src} -> {dst}")
        except Exception as exc:  # Keep batch running even if one file fails.
            fail_count += 1
            print(f"[FAIL] {src}: {exc}")

    return success_count, fail_count


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Batch convert files in data folder to markdown using MarkItDown."
    )
    parser.add_argument(
        "--data-dir",
        default="data",
        type=Path,
        help="Input folder that contains files to convert (default: data)",
    )
    parser.add_argument(
        "--output-dir",
        default="cleaned data",
        type=Path,
        help="Output folder where markdown files are written (default: cleaned data)",
    )
    args = parser.parse_args()

    data_dir = args.data_dir
    output_dir = args.output_dir

    if not data_dir.exists() or not data_dir.is_dir():
        print(f"Input folder not found: {data_dir}")
        return 1

    output_dir.mkdir(parents=True, exist_ok=True)

    success, failed = process_all(data_dir, output_dir)

    print("\nDone.")
    print(f"Converted: {success}")
    print(f"Failed:    {failed}")

    # Non-zero exit when any conversion failed.
    return 0 if failed == 0 else 2


if __name__ == "__main__":
    raise SystemExit(main())
