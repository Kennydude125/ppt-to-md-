from pptx import Presentation
from batch_markitdown import get_smartart_elements, get_text_frame_text, _normalize_text_block, _is_smartart_shape
import json

prs = Presentation("data/3.2 港股非自動對盤交易 Flow v2.1.pptx")
slide = prs.slides[7]
elements_on_slide = []
def process_shape(shape, group_left=None, group_top=None):
    left = group_left if group_left is not None else getattr(shape, "left", 0)
    if getattr(shape, "shape_type", None) == 6:
        g_left = getattr(shape, "left", 0)
        g_top = getattr(shape, "top", 0)
        for i, sub_shape in enumerate(shape.shapes):
            process_shape(sub_shape, group_left=g_left, group_top=g_top + (i * 1000))
        return
    if getattr(shape, "has_text_frame", False):
        text = get_text_frame_text(shape)
        if text:
            elements_on_slide.append({"text": text.replace(chr(10), " ")[:40], "top": group_top if group_top is not None else getattr(shape, "top", 0), "left": left})
        return
    if _is_smartart_shape(shape):
        for sa_item in get_smartart_elements(shape):
            elements_on_slide.append({"text": sa_item["text"].replace(chr(10), " ")[:40], "top": sa_item["top"], "left": sa_item["left"]})

for s in slide.shapes: process_shape(s)
for e in elements_on_slide:
    print(f"{e['left']:9.0f} {e['top']:9.0f} | {e['text']}")
